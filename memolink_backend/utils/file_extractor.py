import io
import re
import html as _html
import base64
import zipfile
from html.parser import HTMLParser

AUDIO_VIDEO_EXTS = {"mp3", "mp4", "mpeg", "mpga", "m4a", "mp4a", "wav", "webm", "ogg", "flac", "avi"}

_WHISPER_MIME: dict[str, str] = {
    "mp3": "audio/mpeg",
    "mp4": "video/mp4",
    "mpeg": "audio/mpeg",
    "mpga": "audio/mpeg",
    "m4a": "audio/mp4",
    "mp4a": "audio/mp4",
    "wav": "audio/wav",
    "webm": "audio/webm",
    "ogg": "audio/ogg",
    "flac": "audio/flac",
    "avi": "video/x-msvideo",
}

_DEEPGRAM_THRESHOLD = 5 * 1024 * 1024   # ≥ 5 MB → Deepgram
_WHISPER_MAX = 25 * 1024 * 1024         # Whisper hard cap


class _HTMLTextExtractor(HTMLParser):
    def __init__(self):
        super().__init__()
        self._parts: list[str] = []
        self._skip = False

    def handle_starttag(self, tag, attrs):
        if tag.lower() in ("script", "style"):
            self._skip = True

    def handle_endtag(self, tag):
        if tag.lower() in ("script", "style"):
            self._skip = False

    def handle_data(self, data):
        if not self._skip:
            text = data.strip()
            if text:
                self._parts.append(text)


class _HTMLSanitizer(HTMLParser):
    """Strip layout tags, keep only TipTap-compatible semantic tags.
    Block layout tags (div, section, etc.) are converted to <p> wrappers
    so their text content is never orphaned."""
    _ALLOWED = {
        "h1", "h2", "h3", "h4", "h5", "h6",
        "p", "br", "hr",
        "ul", "ol", "li",
        "table", "thead", "tbody", "tfoot", "tr", "th", "td",
        "strong", "b", "em", "i", "u", "s", "strike", "code", "pre",
        "blockquote", "a", "img",
    }
    _VOID = {"br", "hr", "img"}
    _SKIP = {"script", "style", "noscript", "svg", "canvas", "select", "option", "nav", "footer"}
    # Layout tags whose text content gets wrapped in <p>
    _BLOCK_LAYOUT = {
        "div", "section", "article", "main", "aside",
        "header", "figure", "figcaption", "details", "summary",
    }
    _SAFE_ATTRS = {
        "a": {"href", "title"},
        "img": {"src", "alt", "width", "height"},
        "td": {"colspan", "rowspan"},
        "th": {"colspan", "rowspan"},
    }
    # Tags that are already block containers — don't wrap inside these
    _BLOCK_CONTAINERS = {
        "h1", "h2", "h3", "h4", "h5", "h6",
        "p", "li", "pre", "blockquote", "td", "th",
    }

    def __init__(self):
        super().__init__()
        self._out: list[str] = []
        self._skip_depth = 0
        self._tag_stack: list[str] = []  # stack of currently open allowed/layout tags

    def _in_block_container(self) -> bool:
        return any(t in self._BLOCK_CONTAINERS for t in self._tag_stack)

    def handle_starttag(self, tag, attrs):
        tag = tag.lower()
        if self._skip_depth > 0:
            if tag in self._SKIP:
                self._skip_depth += 1
            return
        if tag in self._SKIP:
            self._skip_depth += 1
            return
        if tag in self._ALLOWED:
            safe_names = self._SAFE_ATTRS.get(tag, set())
            attr_str = "".join(
                f' {n}="{_html.escape(v or "")}"'
                for n, v in attrs if n in safe_names and v
            )
            self._out.append(f"<{tag}{attr_str}>")
            self._tag_stack.append(tag)
        elif tag in self._BLOCK_LAYOUT:
            if not self._in_block_container():
                self._out.append("<p>")
            self._tag_stack.append(tag)

    def handle_endtag(self, tag):
        tag = tag.lower()
        if self._skip_depth > 0:
            if tag in self._SKIP:
                self._skip_depth -= 1
            return
        if tag in self._ALLOWED and tag not in self._VOID:
            self._out.append(f"</{tag}>")
            if tag in self._tag_stack:
                self._tag_stack.pop()
        elif tag in self._BLOCK_LAYOUT:
            if tag in self._tag_stack:
                self._tag_stack.pop()
            if not self._in_block_container():
                self._out.append("</p>")

    def handle_data(self, data):
        if self._skip_depth > 0:
            return
        self._out.append(_html.escape(data))

    def handle_entityref(self, name):
        if self._skip_depth == 0:
            self._out.append(f"&{name};")

    def handle_charref(self, name):
        if self._skip_depth == 0:
            self._out.append(f"&#{name};")

    def handle_charref(self, name):
        if self._skip_depth == 0:
            self._out.append(f"&#{name};")

    def get_html(self) -> str:
        return "".join(self._out)


def extract_text_local(file_bytes: bytes, filename: str) -> str:
    ext = filename.lower().rsplit(".", 1)[-1]

    if ext in ("txt", "md"):
        return file_bytes.decode("utf-8", errors="ignore")

    if ext in ("html", "htm"):
        try:
            parser = _HTMLTextExtractor()
            parser.feed(file_bytes.decode("utf-8", errors="ignore"))
            return "\n".join(parser._parts)
        except Exception as e:
            return f"[HTML extraction error] {e}"

    if ext == "pdf":
        try:
            import pdfplumber
            parts = []
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                for page in pdf.pages:
                    parts.append(page.extract_text() or "")
            return "\n".join(parts)
        except Exception as e:
            return f"[PDF extraction error] {e}"

    if ext == "docx":
        try:
            import docx
            doc = docx.Document(io.BytesIO(file_bytes))
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception as e:
            return f"[DOCX extraction error] {e}"

    if ext == "pptx":
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        parts.append(shape.text)
            return "\n".join(parts)
        except Exception as e:
            return f"[PPTX extraction error] {e}"

    if ext == "ppt":
        return (
            "Old .ppt format is not supported directly. "
            "Please open in PowerPoint, save as .pptx, and re-upload."
        )

    if ext == "zip":
        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                return "\n".join(
                    z.read(name).decode("utf-8", errors="ignore")
                    for name in z.namelist()
                    if name.lower().endswith(".txt")
                )
        except Exception as e:
            return f"[ZIP extraction error] {e}"

    if ext in AUDIO_VIDEO_EXTS:
        result = transcribe_audio(file_bytes, filename, ext)
        if result and not result.startswith("["):
            return f"[Audio transcription: {filename}]\n\n{result}"
        return result

    return ""


# ---------------------------------------------------------------------------
# Rich HTML extraction — preserves headings, bold, italic, lists, tables
# ---------------------------------------------------------------------------

def _plain_to_html(text: str) -> str:
    lines = text.splitlines()
    return "".join(
        f"<p>{_html.escape(line)}</p>" if line.strip() else "<p></p>"
        for line in lines
    ) or f"<p>{_html.escape(text)}</p>"


def _runs_to_html(runs) -> str:
    result = ""
    for run in runs:
        text = _html.escape(run.text)
        if not text:
            continue
        if run.bold:
            text = f"<strong>{text}</strong>"
        if run.italic:
            text = f"<em>{text}</em>"
        if run.underline:
            text = f"<u>{text}</u>"
        result += text
    return result


def _para_to_html(p) -> str:
    style = (p.style.name or "").lower() if p.style else ""
    inline = _runs_to_html(p.runs)
    if not inline.strip():
        return "<p></p>"
    if "heading 1" in style:
        return f"<h1>{inline}</h1>"
    if "heading 2" in style:
        return f"<h2>{inline}</h2>"
    if "heading 3" in style:
        return f"<h3>{inline}</h3>"
    if style.startswith("heading"):
        return f"<h4>{inline}</h4>"
    if "list bullet" in style or "list paragraph" in style:
        return f"<ul><li>{inline}</li></ul>"
    if "list number" in style:
        return f"<ol><li>{inline}</li></ol>"
    return f"<p>{inline}</p>"


def _table_to_html(tbl) -> str:
    rows = []
    for i, row in enumerate(tbl.rows):
        cells = []
        for cell in row.cells:
            text = _html.escape(" ".join(p.text for p in cell.paragraphs).strip())
            tag = "th" if i == 0 else "td"
            cells.append(f"<{tag}>{text}</{tag}>")
        rows.append(f"<tr>{''.join(cells)}</tr>")
    return f"<table>{''.join(rows)}</table>"


_BLIP_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
_REL_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


def _para_images_html(para, doc) -> list[str]:
    """Return <img> tags for any inline images embedded in this DOCX paragraph."""
    imgs = []
    for blip in para._element.iter(f"{_BLIP_NS}blip"):
        embed = blip.get(f"{_REL_NS}embed")
        if not embed:
            continue
        try:
            img_part = doc.part.related_parts[embed]
            b64 = base64.b64encode(img_part.blob).decode()
            ct = img_part.content_type
            imgs.append(f'<img src="data:{ct};base64,{b64}" style="max-width:100%">')
        except Exception:
            pass
    return imgs


def _docx_to_html(file_bytes: bytes) -> str:
    import docx  # python-docx
    doc = docx.Document(io.BytesIO(file_bytes))
    parts = []
    para_iter = iter(doc.paragraphs)
    table_iter = iter(doc.tables)
    for child in doc.element.body:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag == "p":
            try:
                para = next(para_iter)
                parts.extend(_para_images_html(para, doc))
                parts.append(_para_to_html(para))
            except StopIteration:
                pass
        elif tag == "tbl":
            try:
                parts.append(_table_to_html(next(table_iter)))
            except StopIteration:
                pass
    html = "".join(parts)
    html = re.sub(r"</ul>\s*<ul>", "", html)
    html = re.sub(r"</ol>\s*<ol>", "", html)
    return html


def _pdf_to_html(file_bytes: bytes) -> str:
    import pdfplumber
    parts = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            words = page.extract_words(extra_attrs=["fontname", "size"])
            if not words:
                for line in (page.extract_text() or "").splitlines():
                    if line.strip():
                        parts.append(f"<p>{_html.escape(line.strip())}</p>")
                continue

            sizes = [float(w.get("size") or 0) for w in words if w.get("size")]
            body_size = sorted(sizes)[len(sizes) // 2] if sizes else 12.0

            line_map: dict[float, list] = {}
            for w in words:
                top = round(float(w.get("top") or 0), 0)
                line_map.setdefault(top, []).append(w)

            for top in sorted(line_map):
                line_words = sorted(line_map[top], key=lambda w: float(w.get("x0") or 0))
                text = " ".join(w["text"] for w in line_words).strip()
                if not text:
                    continue
                avg_size = sum(float(w.get("size") or body_size) for w in line_words) / len(line_words)
                is_bold = any(
                    "bold" in (w.get("fontname") or "").lower() for w in line_words
                )
                escaped = _html.escape(text)
                if avg_size >= body_size * 1.5:
                    parts.append(f"<h1>{escaped}</h1>")
                elif avg_size >= body_size * 1.3:
                    parts.append(f"<h2>{escaped}</h2>")
                elif avg_size >= body_size * 1.15 or (is_bold and len(text) < 80):
                    parts.append(f"<h3>{escaped}</h3>")
                elif is_bold:
                    parts.append(f"<p><strong>{escaped}</strong></p>")
                else:
                    parts.append(f"<p>{escaped}</p>")
    return "".join(parts)


def _pptx_to_html(file_bytes: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    total = len(prs.slides)
    for slide_num, slide in enumerate(prs.slides, 1):
        title_text = None
        content_items: list[tuple[int, str]] = []
        slide_images: list[str] = []
        for shape in slide.shapes:
            # Extract embedded images from picture shapes
            if hasattr(shape, "image"):
                try:
                    img_data = shape.image.blob
                    ext = shape.image.ext.lower()
                    mime = "image/jpeg" if ext in ("jpg", "jpeg") else f"image/{ext}"
                    b64 = base64.b64encode(img_data).decode()
                    slide_images.append(f'<img src="data:{mime};base64,{b64}" style="max-width:100%">')
                except Exception:
                    pass
                continue
            if not hasattr(shape, "text_frame"):
                continue
            ph = getattr(shape, "placeholder_format", None)
            if ph is not None and ph.idx == 0:
                title_text = shape.text.strip()
            else:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        content_items.append((para.level, text))
        parts.append(f"<h2>{_html.escape(title_text or f'Slide {slide_num}')}</h2>")
        if content_items:
            items_html = "".join(
                f"<li>{'&nbsp;&nbsp;' * level}{_html.escape(text)}</li>"
                for level, text in content_items
            )
            parts.append(f"<ul>{items_html}</ul>")
        parts.extend(slide_images)
        if slide_num < total:
            parts.append("<hr>")
    return "".join(parts)


def extract_formatted_html(file_bytes: bytes, filename: str) -> str:
    """
    Extract document content as rich HTML for storage in the note editor.
    Preserves headings, bold/italic, bullet lists, numbered lists, and tables.
    Falls back to plain-text paragraph wrapping for unsupported types.
    """
    ext = filename.lower().rsplit(".", 1)[-1]

    if ext == "docx":
        try:
            return _docx_to_html(file_bytes)
        except Exception:
            pass

    if ext == "pdf":
        try:
            return _pdf_to_html(file_bytes)
        except Exception:
            pass

    if ext == "pptx":
        try:
            return _pptx_to_html(file_bytes)
        except Exception:
            pass

    if ext in ("html", "htm"):
        try:
            import html2text
            raw = file_bytes.decode("utf-8", errors="ignore")
            converter = html2text.HTML2Text()
            converter.ignore_links = False
            converter.ignore_images = False
            converter.body_width = 0
            return converter.handle(raw).strip()
        except Exception:
            pass

    # For all other types (txt, md, zip, audio) use plain extraction + wrap
    text = extract_text_local(file_bytes, filename)
    return _plain_to_html(text)


_DEEPGRAM_MIME: dict[str, str] = {
    "mp3": "audio/mpeg", "mp4": "video/mp4", "m4a": "audio/mp4",
    "wav": "audio/wav", "webm": "audio/webm", "ogg": "audio/ogg",
    "flac": "audio/flac", "mpeg": "audio/mpeg", "mpga": "audio/mpeg",
    "mov": "video/quicktime", "avi": "video/x-msvideo",
}


def _transcribe_deepgram(file_bytes: bytes, filename: str, ext: str, language: str | None = None) -> str:
    import httpx
    from memolink_backend.core.config import settings

    content_type = _DEEPGRAM_MIME.get(ext.lstrip("."), "audio/mpeg")
    params: dict = {"model": "nova-2", "smart_format": "true"}
    if language:
        params["language"] = language

    resp = httpx.post(
        "https://api.deepgram.com/v1/listen",
        params=params,
        content=file_bytes,
        headers={
            "Authorization": f"Token {settings.deepgram_api_key}",
            "Content-Type": content_type,
        },
        timeout=300.0,
    )
    resp.raise_for_status()
    data = resp.json()
    return data["results"]["channels"][0]["alternatives"][0]["transcript"]


def _transcribe_whisper(file_bytes: bytes, filename: str, language: str | None = None) -> str:
    from openai import OpenAI
    from memolink_backend.core.config import settings

    client = OpenAI(api_key=settings.openai_api_key)
    safe_filename = re.sub(r"\s+", "_", filename)
    buf = io.BytesIO(file_bytes)
    buf.name = safe_filename
    kwargs: dict = {"model": "whisper-1", "file": buf}
    if language:
        kwargs["language"] = language
    transcript = client.audio.transcriptions.create(**kwargs)
    return transcript.text


def transcribe_audio(file_bytes: bytes, filename: str, ext: str, language: str | None = None) -> str:
    """< 5 MB → Whisper; ≥ 5 MB → Deepgram (fallback: Whisper if ≤ 25 MB)."""
    from memolink_backend.core.config import settings

    size = len(file_bytes)

    if size < _DEEPGRAM_THRESHOLD:
        try:
            return _transcribe_whisper(file_bytes, filename, language)
        except Exception as e:
            detail = getattr(e, "body", None) or ""
            return f"[Transcription error] {e}{f' — detail: {detail}' if detail else ''}"

    # ≥ 5 MB: try Deepgram first
    if settings.deepgram_api_key:
        try:
            return _transcribe_deepgram(file_bytes, filename, ext, language)
        except Exception:
            pass  # fall through to Whisper fallback

    # Whisper fallback (only if within its 25 MB cap)
    if size <= _WHISPER_MAX:
        try:
            return _transcribe_whisper(file_bytes, filename, language)
        except Exception as e:
            detail = getattr(e, "body", None) or ""
            return f"[Transcription error] {e}{f' — detail: {detail}' if detail else ''}"

    return "[Transcription skipped] File exceeds the 25 MB limit and Deepgram is not configured or failed."
